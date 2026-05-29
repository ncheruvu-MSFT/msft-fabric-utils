"""Build a Microsoft-styled PPT for the Fabric Target Architecture summary.

Visual style follows the attached sample template:
- Dark navy background (#2B3A4A)
- Teal accent (#7FEFD6) for titles
- White body text, light-gray sub text
- Microsoft wordmark on title slide
- Rounded "card" callouts with teal left bars
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Theme ----------
BG        = RGBColor(0x2B, 0x3A, 0x4A)
ACCENT    = RGBColor(0x7F, 0xEF, 0xD6)
ACCENT2   = RGBColor(0x4F, 0xC3, 0xA1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUBTEXT   = RGBColor(0xC9, 0xD3, 0xDC)
CARD_BG   = RGBColor(0x33, 0x44, 0x55)
CARD_LINE = RGBColor(0x55, 0x6B, 0x7E)
PHASE_RED = RGBColor(0xE0, 0x4B, 0x3A)
PHASE_ORG = RGBColor(0xEA, 0x7A, 0x2B)
PHASE_BLU = RGBColor(0x2E, 0x4C, 0xB8)
PHASE_LBL = RGBColor(0x2A, 0x9D, 0xE0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def make_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=WHITE, bullet_char="\u25B8", bullet_color=ACCENT,
                line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet_char}  "
        r.font.name = "Segoe UI"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height, *,
             fill=CARD_BG, line=CARD_LINE, accent_bar=ACCENT, bar_side="left"):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(0.75)
    if accent_bar is not None:
        bar_w = Emu(50000)
        if bar_side == "left":
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Emu(30000),
                                         bar_w, height - Emu(60000))
        else:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         left + width - bar_w, top + Emu(30000),
                                         bar_w, height - Emu(60000))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_bar
        bar.line.fill.background()
    return card


def add_circle(slide, cx, cy, d, fill, text):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"


def add_page_number(slide, n):
    add_text(slide, Inches(12.7), Inches(7.15), Inches(0.5), Inches(0.3),
             str(n), size=10, color=SUBTEXT, align=PP_ALIGN.RIGHT)


def title(slide, txt, top=Inches(0.45)):
    add_text(slide, Inches(0.55), top, Inches(12.2), Inches(0.9),
             txt, size=40, bold=True, color=ACCENT)


# ---------- Slides ----------
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text(s, Inches(0.55), Inches(0.45), Inches(3), Inches(0.4),
             "Microsoft", size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.55), Inches(2.0), Inches(12), Inches(2.8),
             "Fabric Target\nArchitecture",
             size=60, bold=True, color=ACCENT)
    add_text(s, Inches(0.55), Inches(5.0), Inches(12), Inches(1.0),
             "OneLake EDW \u00B7 Medallion Pipelines \u00B7 Shortcut-Based Consumption \u00B7 Governed by Purview",
             size=18, color=SUBTEXT)
    add_text(s, Inches(0.55), Inches(5.6), Inches(12), Inches(0.6),
             "A Phased Path to a Reliable, Cost-Efficient, Governed Data Foundation",
             size=14, color=SUBTEXT)
    add_page_number(s, 1)


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.4), Inches(0.5), Inches(4.0), Inches(6.5))
    panel.adjustments[0] = 0.06
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x22, 0x2F, 0x3D)
    panel.line.color.rgb = CARD_LINE
    add_text(s, Inches(0.7), Inches(3.0), Inches(3.5), Inches(1.2),
             "Agenda", size=54, bold=True, color=ACCENT)

    items = [
        ("The Challenge", "Why the current state is hard to scale, govern, and operate cost-effectively"),
        ("Proposed Target Architecture", "Centralized OneLake EDW + medallion pipelines + shortcut consumption"),
        ("Pipeline & Semantic Model Design", "Separation of stages and Copilot-assisted optimization"),
        ("Access, Workspaces & Governance", "Workspace strategy, RBAC, Fabric catalog, Purview layering"),
        ("Phased Implementation Plan", "Four phases from foundation to enterprise catalog expansion"),
        ("Expected Outcomes & Next Steps", "Measurable reliability, cost, and time-to-analytics wins"),
    ]
    x = Inches(4.9); y = Inches(0.8)
    for head, sub in items:
        add_text(s, x, y, Inches(8), Inches(0.45), f"\u2022  {head}",
                 size=22, bold=True, color=WHITE)
        add_text(s, x + Inches(0.4), y + Inches(0.5), Inches(8), Inches(0.4),
                 sub, size=13, color=SUBTEXT)
        y += Inches(1.0)
    add_page_number(s, 2)


def slide_challenge(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "The Challenge")
    add_bullets(s, Inches(0.6), Inches(1.6), Inches(12), Inches(4.5), [
        "Data sprawl across multiple workspaces and copies \u2014 storage and compute duplication",
        "Inconsistent data quality and schemas; type fixes deferred to semantic models",
        "\u201CTriple-write\u201D anti-pattern in pipelines \u2014 unclear ownership, harder debugging",
        "Semantic models doing heavy lifting (e.g., VARCHAR-heavy) instead of upstream fixes",
        "No unified discovery and governance across Fabric and other clouds (AWS, M365, etc.)",
        "Cost and performance pressure without a clear, governed architectural baseline",
    ], size=20)

    add_card(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.3),
             fill=RGBColor(0x32, 0x44, 0x55), accent_bar=ACCENT)
    add_text(s, Inches(0.95), Inches(6.0), Inches(11.5), Inches(0.35),
             "CORE QUESTION", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.95), Inches(6.35), Inches(11.5), Inches(0.9),
             "How do we centralize the data plane, fix quality upstream, and govern access \u2014 "
             "without re-platforming or over-engineering?",
             size=15, color=WHITE)
    add_page_number(s, 3)


def slide_approach(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Our Approach")

    cards = [
        ("CENTRAL DATA PLANE",       "Single OneLake-based EDW (one copy of truth)"),
        ("DATA QUALITY",             "Medallion: Bronze \u2192 Silver \u2192 Gold inside OneLake"),
        ("CONSUMPTION",              "Shortcuts (not copies) to downstream workspaces"),
        ("OPTIMIZATION",             "Copilot-assisted semantic model tuning"),
        ("ACCESS & RBAC",            "Workspace-scoped access, dataset-level RBAC"),
        ("GOVERNANCE LAYER",         "Fabric catalog first \u2192 Purview for cross-cloud"),
    ]
    col_w = Inches(6.05); row_h = Inches(1.1)
    for i, (head, body) in enumerate(cards):
        r, c = divmod(i, 2)
        left = Inches(0.55) + (col_w + Inches(0.15)) * c
        top  = Inches(1.5) + (row_h + Inches(0.2)) * r
        add_card(s, left, top, col_w, row_h, accent_bar=ACCENT)
        add_text(s, left + Inches(0.35), top + Inches(0.15),
                 col_w - Inches(0.5), Inches(0.35),
                 head, size=11, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.35), top + Inches(0.48),
                 col_w - Inches(0.5), Inches(0.55),
                 body, size=15, color=WHITE)

    add_card(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(1.2),
             fill=RGBColor(0x32, 0x44, 0x55), accent_bar=ACCENT)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.35),
             "KEY INSIGHT", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.8),
             "Fix data quality and typing upstream in Silver/Gold. Semantic models accelerate \u2014 "
             "they don't compensate for weak foundations.",
             size=14, color=WHITE)
    add_page_number(s, 4)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Target Architecture \u2014 OneLake + Medallion")

    add_card(s, Inches(0.4), Inches(1.7), Inches(2.4), Inches(4.6),
             fill=RGBColor(0x2E, 0x44, 0x60), accent_bar=ACCENT2)
    add_text(s, Inches(0.55), Inches(1.85), Inches(2.2), Inches(0.4),
             "SOURCES", size=12, bold=True, color=ACCENT)
    for i, src in enumerate(["Azure", "AWS", "Workday", "CRM", "SaaS / Files", "Streaming"]):
        add_text(s, Inches(0.7), Inches(2.35) + Inches(0.55) * i,
                 Inches(2.0), Inches(0.4), f"\u2022 {src}", size=14, color=WHITE)

    add_card(s, Inches(3.2), Inches(1.7), Inches(6.7), Inches(4.6),
             fill=RGBColor(0x22, 0x3A, 0x52), accent_bar=ACCENT, bar_side="left")
    add_text(s, Inches(3.4), Inches(1.85), Inches(6.3), Inches(0.45),
             "ONELAKE \u2014 ENTERPRISE DATA PLANE (EDW)",
             size=13, bold=True, color=ACCENT)

    layers = [
        ("BRONZE", "Raw ingested data\nfrom source systems", RGBColor(0x8B, 0x5A, 0x2B)),
        ("SILVER", "Cleansed, standardized\nschema + types enforced", RGBColor(0x9A, 0xA5, 0xB1)),
        ("GOLD",   "Analytics-ready,\nbusiness-curated datasets", RGBColor(0xD4, 0xAF, 0x37)),
    ]
    pill_w = Inches(2.05); pill_h = Inches(2.9)
    for i, (name, body, col) in enumerate(layers):
        left = Inches(3.35) + (pill_w + Inches(0.1)) * i
        top  = Inches(2.45)
        add_card(s, left, top, pill_w, pill_h, fill=RGBColor(0x2D, 0x3E, 0x52),
                 accent_bar=col)
        add_text(s, left, top + Inches(0.25), pill_w, Inches(0.5),
                 name, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), top + Inches(0.9),
                 pill_w - Inches(0.4), Inches(1.6),
                 body, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(3.35), Inches(5.55), Inches(4.2), Inches(0.35),
             "Pipeline A: Bronze \u2192 Silver", size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.55), Inches(5.85), Inches(4.2), Inches(0.35),
             "Pipeline B: Silver \u2192 Gold", size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_card(s, Inches(10.25), Inches(1.7), Inches(2.65), Inches(4.6),
             fill=RGBColor(0x2E, 0x44, 0x60), accent_bar=ACCENT2, bar_side="right")
    add_text(s, Inches(10.4), Inches(1.85), Inches(2.4), Inches(0.4),
             "CONSUMERS", size=12, bold=True, color=ACCENT)
    consumers = ["Power BI / Semantic", "Analytics teams", "Data science / ML",
                 "Business workspaces", "Shortcuts \u2014 no copies"]
    for i, c in enumerate(consumers):
        add_text(s, Inches(10.45), Inches(2.35) + Inches(0.55) * i,
                 Inches(2.4), Inches(0.45), f"\u2022 {c}", size=13, color=WHITE)

    add_card(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.7),
             fill=RGBColor(0x32, 0x44, 0x55), accent_bar=ACCENT)
    add_text(s, Inches(0.75), Inches(6.62), Inches(12.0), Inches(0.6),
             "Governance: Fabric Catalog & Lineage (native)   \u2192   Purview (cross-cloud, glossary, certifications)",
             size=13, bold=True, color=WHITE)
    add_page_number(s, 5)


def slide_pipeline_pattern(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Pipeline Design Pattern")

    for i, (name, body, items) in enumerate([
        ("Pipeline A \u2014 Bronze \u2192 Silver",
         "Ingestion & standardization",
         ["Source-system ingestion", "Schema enforcement & type fixes",
          "Cleansing and de-duplication", "Quality checks at the boundary"]),
        ("Pipeline B \u2014 Silver \u2192 Gold",
         "Business logic & analytics readiness",
         ["Joins, aggregations, business rules", "Conformed dimensions / metrics",
          "Optimized for semantic models & ML", "Owned by analytics engineering"]),
    ]):
        left = Inches(0.55) + (Inches(6.15) + Inches(0.15)) * i
        add_card(s, left, Inches(1.5), Inches(6.15), Inches(3.6),
                 accent_bar=ACCENT)
        add_text(s, left + Inches(0.35), Inches(1.65), Inches(5.6), Inches(0.5),
                 name, size=20, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.35), Inches(2.15), Inches(5.6), Inches(0.4),
                 body, size=13, color=SUBTEXT)
        add_bullets(s, left + Inches(0.35), Inches(2.7), Inches(5.6), Inches(2.3),
                    items, size=14)

    add_card(s, Inches(0.55), Inches(5.3), Inches(12.3), Inches(1.85),
             fill=RGBColor(0x46, 0x33, 0x33), accent_bar=PHASE_RED)
    add_text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.4),
             "AVOID \u2014 \u201CTriple-Write\u201D Anti-Pattern", size=13, bold=True, color=PHASE_RED)
    add_bullets(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.3), [
        "Single pipelines writing Bronze/Silver/Gold together \u2192 unclear ownership, harder debugging",
        "Mixing cleansing with business logic blocks safe long-term maintenance and rollback",
    ], size=13)
    add_page_number(s, 6)


def slide_semantic(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Semantic Model Optimization \u2014 Acceleration, Not a Crutch")

    add_bullets(s, Inches(0.6), Inches(1.6), Inches(12), Inches(3.5), [
        "Use Fabric\u2019s Copilot-assisted semantic model optimization to surface common issues",
        "Typical findings: VARCHAR-heavy columns, missing data types, weak relationships",
        "Push schema fixes back into Silver/Gold \u2014 with a human review step",
        "Optimization is a remediation accelerator, not a replacement for upstream quality",
    ], size=18)

    for i, (head, body, color) in enumerate([
        ("Anti-Pattern",
         "Mask upstream issues inside the semantic model \u2192\nshort-term wins, long-term tech debt",
         PHASE_RED),
        ("Recommended",
         "Use Copilot to detect \u2192 fix in Silver/Gold \u2192\nsemantic model stays lean and fast",
         ACCENT2),
    ]):
        left = Inches(0.55) + (Inches(6.15) + Inches(0.15)) * i
        add_card(s, left, Inches(5.0), Inches(6.15), Inches(2.0),
                 accent_bar=color)
        add_text(s, left + Inches(0.35), Inches(5.15), Inches(5.6), Inches(0.45),
                 head, size=16, bold=True, color=color)
        add_text(s, left + Inches(0.35), Inches(5.65), Inches(5.6), Inches(1.3),
                 body, size=14, color=WHITE)
    add_page_number(s, 7)


def slide_access_governance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Access, Workspaces & Governance")

    add_text(s, Inches(0.6), Inches(1.5), Inches(6), Inches(0.5),
             "Workspaces & Access", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(2.05), Inches(6), Inches(3.5), [
        "Central data engineering owns Bronze/Silver/Gold",
        "Business / analytics teams in separate Fabric workspaces",
        "Consume via OneLake shortcuts \u2014 no physical copies",
        "RBAC enforced at workspace and dataset level",
        "Data remains centrally governed and observable",
    ], size=16)

    add_text(s, Inches(7.0), Inches(1.5), Inches(6), Inches(0.5),
             "Governance Layering", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.0), Inches(2.05), Inches(6), Inches(3.5), [
        "Fabric native catalog & lineage = primary visibility",
        "Layer Purview after the foundation is sound",
        "Purview adds cross-platform reach (AWS, Azure, M365)",
        "Business glossary, certifications, enterprise discovery",
        "Avoid bolting governance onto unstable data",
    ], size=16)

    add_card(s, Inches(0.55), Inches(6.0), Inches(12.3), Inches(1.2),
             fill=RGBColor(0x32, 0x44, 0x55), accent_bar=ACCENT)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.35),
             "PRINCIPLE", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.8),
             "Govern the data plane once, distribute access many times \u2014 with Fabric first and Purview as the enterprise overlay.",
             size=14, color=WHITE)
    add_page_number(s, 8)


def slide_phases(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Phased Implementation Plan")

    phases = [
        ("1", "Foundation & Architecture Alignment",
         "Confirm OneLake EDW + medallion. Define workspace strategy. "
         "Identify priority pipelines/datasets to remediate first.", PHASE_RED),
        ("2", "Data Quality & Pipeline Refactor",
         "Refactor pipelines one domain at a time. Enforce schemas/types in Silver. "
         "Eliminate triple-writes. Apply semantic model optimization selectively.", PHASE_ORG),
        ("3", "Consumption & Governance Enablement",
         "Migrate analytics teams to shortcut-based consumption. Apply consistent RBAC. "
         "Enable Fabric catalog and lineage for transparency.", PHASE_BLU),
        ("4", "Enterprise Catalog Expansion",
         "Integrate Purview for cross-cloud and non-Fabric sources. "
         "Curate certified datasets, business definitions, and ownership.", PHASE_LBL),
    ]
    col_w = Inches(3.05); col_h = Inches(5.0)
    for i, (num, head, body, color) in enumerate(phases):
        left = Inches(0.4) + (col_w + Inches(0.12)) * i
        top  = Inches(1.6)
        add_card(s, left, top, col_w, col_h, accent_bar=color)
        add_circle(s, left + Inches(0.6), top + Inches(0.55),
                   Inches(0.7), color, num)
        add_text(s, left + Inches(0.3), top + Inches(1.2),
                 col_w - Inches(0.5), Inches(1.2),
                 head, size=16, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.3), top + Inches(2.5),
                 col_w - Inches(0.5), Inches(2.4),
                 body, size=13, color=WHITE)
    add_page_number(s, 9)


def slide_outcomes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Expected Outcomes")

    chips = ["Reliable Data", "Lower Storage Cost", "Lower Compute Cost",
             "Faster Time-to-Analytics", "Clear Ownership", "Right-Sized Governance"]
    cw = Inches(2.9); ch = Inches(0.5)
    for i, c in enumerate(chips):
        r, col = divmod(i, 4)
        left = Inches(0.55) + (cw + Inches(0.2)) * col
        top  = Inches(1.5) + (ch + Inches(0.15)) * r
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cw, ch)
        chip.adjustments[0] = 0.4
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(0x2E, 0x44, 0x55)
        chip.line.color.rgb = ACCENT
        tf = chip.text_frame; tf.margin_left=tf.margin_right=Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run(); r1.text = f"\u2713  {c}"
        r1.font.name="Segoe UI"; r1.font.size=Pt(13); r1.font.color.rgb = ACCENT

    add_text(s, Inches(0.6), Inches(3.3), Inches(6), Inches(0.5),
             "Business Impact", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(0.6), Inches(3.8), Inches(6), Inches(3.0), [
        "Improved data reliability and performance",
        "Less semantic-model \u201Cheavy lifting\u201D at query time",
        "Faster delivery of trusted reports and ML features",
        "Single source of truth for cross-team analytics",
    ], size=15)

    add_text(s, Inches(7.0), Inches(3.3), Inches(6), Inches(0.5),
             "Operational Impact", size=20, bold=True, color=ACCENT)
    add_bullets(s, Inches(7.0), Inches(3.8), Inches(6), Inches(3.0), [
        "Reduced storage & compute via shortcut-based consumption",
        "Clear pipeline ownership, easier debugging and rollback",
        "Right-sized governance \u2014 Fabric first, Purview enterprise-wide",
        "Lower risk of bolting governance onto unstable data",
    ], size=15)
    add_page_number(s, 10)


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title(s, "Summary & Next Steps")
    add_bullets(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(4.0), [
        "Centralize the data plane in OneLake \u2014 one copy, governed centrally",
        "Enforce quality upstream with a clean Bronze/Silver/Gold medallion",
        "Split pipelines by stage \u2014 avoid the triple-write anti-pattern",
        "Use Copilot semantic model optimization to accelerate remediation, not mask issues",
        "Distribute access through shortcuts and RBAC \u2014 not data copies",
        "Layer Purview after the foundation is sound for cross-cloud discovery",
    ], size=18)

    add_card(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.2),
             fill=RGBColor(0x32, 0x44, 0x55), accent_bar=ACCENT)
    add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.35),
             "NEXT STEPS", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(6.02), Inches(11.5), Inches(0.8),
             "Confirm target architecture, pick a pilot domain for Phase 2 pipeline refactor, "
             "and align on the workspace + RBAC model before Phase 3.",
             size=14, color=WHITE)

    add_text(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.4),
             "Q&A", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_page_number(s, 11)


def build(out_path: str):
    prs = make_prs()
    slide_title(prs)
    slide_agenda(prs)
    slide_challenge(prs)
    slide_approach(prs)
    slide_architecture(prs)
    slide_pipeline_pattern(prs)
    slide_semantic(prs)
    slide_access_governance(prs)
    slide_phases(prs)
    slide_outcomes(prs)
    slide_summary(prs)
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    import os, sys
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        os.path.dirname(__file__), "..", "..", "docs", "fabric-target-architecture.pptx")
    out = os.path.abspath(out)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    build(out)
