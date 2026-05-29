"""Rebuild fabric-engagement-summary.pptx using the reference deck's master/theme.

Strategy:
  1. Clone the reference deck (preserves all masters, layouts, themes, fonts, images).
  2. Remove every existing slide from the clone.
  3. For each engagement slide (extracted from .bak.pptx into _engagement_content.json),
     add a new slide using a matching reference layout:
       - slide 0  -> 'Title Slide - AI 2'   (master 0, layout 12)
       - others   -> 'Title and Content'    (master 1, layout 15)
  4. Title text = first bullet (heading). Remaining bullets fill body placeholder.
  5. Save over docs/fabric-engagement-summary.pptx.
"""
from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

ROOT = Path(r"c:\Git\AZ\msft-fabric-utils")
REF  = ROOT / "docs" / "Cloud & AI Platforms Envisioning Quick Reference Guide.PPTX"
DST  = ROOT / "docs" / "fabric-engagement-summary.pptx"
CONTENT = ROOT / "scripts" / "python" / "_engagement_content.json"

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def strip_all_slides(prs):
    """Delete every slide from the presentation cleanly."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(prs.slides)
    # remove relationships + parts
    for slide in slides:
        rId = None
        for sid in list(sldIdLst):
            if sid.get(f"{{{NS_R}}}id") and prs.part.related_part(sid.get(f"{{{NS_R}}}id")) is slide.part:
                rId = sid.get(f"{{{NS_R}}}id")
                sldIdLst.remove(sid)
                break
        if rId:
            prs.part.drop_rel(rId)


def find_layout(prs, master_idx, layout_idx):
    return prs.slide_masters[master_idx].slide_layouts[layout_idx]


def find_layout_by_name(prs, name):
    for m in prs.slide_masters:
        for lo in m.slide_layouts:
            if lo.name == name:
                return lo
    raise KeyError(name)


def set_placeholder_text(slide, ph_idx_or_name, lines, *, body_size=None):
    """Find placeholder by idx (int) or name (str) and fill with lines."""
    for ph in slide.placeholders:
        match = (
            (isinstance(ph_idx_or_name, int) and ph.placeholder_format.idx == ph_idx_or_name)
            or (isinstance(ph_idx_or_name, str) and ph.name == ph_idx_or_name)
        )
        if not match:
            continue
        tf = ph.text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            if body_size:
                for r in p.runs:
                    r.font.size = body_size
        return ph
    return None


def first_body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            return ph
    return None


def main():
    if not REF.exists():
        raise SystemExit(f"missing reference: {REF}")
    if not CONTENT.exists():
        raise SystemExit(f"missing content json: {CONTENT}")
    data = json.loads(CONTENT.read_text(encoding="utf-8"))

    # 1. clone reference
    shutil.copy2(REF, DST)
    prs = Presentation(str(DST))

    # 2. strip existing slides
    strip_all_slides(prs)

    title_layout = find_layout_by_name(prs, "Title Slide - AI 2")
    content_layout = find_layout_by_name(prs, "Title and Content")

    for entry in data:
        idx = entry["idx"]
        bullets = entry["bullets"][:]
        title = bullets.pop(0) if bullets else f"Slide {idx+1}"
        # Strip the trailing "Naga Venkata Cheruvu, CSA  ·  Microsoft" attribution
        bullets = [b for b in bullets if "CSA  ·  Microsoft" not in b]

        if idx == 0:
            slide = prs.slides.add_slide(title_layout)
            # title
            set_placeholder_text(slide, 0, ["Fabric Optimization Planning"])
            # subtitle = participants + date
            set_placeholder_text(slide, 12, bullets, body_size=Pt(14))
        else:
            slide = prs.slides.add_slide(content_layout)
            set_placeholder_text(slide, 0, [title])
            body = first_body_placeholder(slide)
            if body is not None:
                tf = body.text_frame
                tf.clear()
                tf.word_wrap = True
                # density-based base size so dense slides shrink automatically
                n = len(bullets)
                if n <= 10:
                    base, header = 14, 16
                elif n <= 16:
                    base, header = 12, 14
                elif n <= 22:
                    base, header = 10, 12
                else:
                    base, header = 9, 10
                for i, line in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    is_section = (
                        not line.startswith("▸")
                        and not line.startswith("•")
                        and not line.startswith("https://")
                        and len(line) < 60
                        and (line.isupper() or line.endswith(":") or line.istitle())
                    )
                    for r in p.runs:
                        r.font.size = Pt(header if is_section else base)
                        if is_section:
                            r.font.bold = True
                # Enable PowerPoint's autofit so any residual overflow shrinks
                bodyPr = tf._txBody.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
                )
                if bodyPr is not None:
                    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
                        for el in bodyPr.findall(f"{{{A}}}{tag}"):
                            bodyPr.remove(el)
                    etree.SubElement(bodyPr, f"{{{A}}}normAutofit",
                                     {"fontScale": "85000", "lnSpcReduction": "10000"})

    prs.save(str(DST))
    print(f"rebuilt: {DST.name}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
